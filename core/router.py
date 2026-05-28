# ============================================
# MANTHAN AI — Multi-Model AI Router
# File: core/router.py
# ============================================

import httpx
import base64
import json
import time
import os
import re
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from .model import get_provider, ProviderConfig, PROVIDERS

async def ask_ai(
    task: str,
    user_message: str,
    context: str = "",
    file_bytes: bytes = None,
    mime_type: str = None
) -> dict:

    # Image hai to Gemini Vision use karo
    if file_bytes and mime_type and mime_type.startswith("image/"):
        provider = PROVIDERS["gemini"]
        return await _call_gemini_vision(provider, user_message, file_bytes, mime_type)

    # Context hai to message me add karo
    full_message = user_message
    if context:
        full_message = f"Context:\n{context}\n\nQuestion:\n{user_message}"

    provider = get_provider(task)

    if task == "ppt":
        if not Presentation:
            return {"response": "PPTX module nahi mili! `pip install python-pptx` run karo."}
        
        system_prompt = (
            "You are an AI that creates PPT slides. You MUST output ONLY raw JSON data in this exact format, with no markdown formatting, no code blocks (like ```json), just the plain JSON array: "
            "[{\"title\": \"Slide 1 Title\", \"content\": [\"Point 1\", \"Point 2\"]}, {\"title\": \"Slide 2 Title\", \"content\": [\"Point 1\"]}]"
        )
        if provider.name == "Google Gemini":
            res = await _call_gemini(provider, system_prompt, f"Create a presentation on: {user_message}. Output plain JSON only.")
        else:
            res = await _call_openai_compatible(provider, system_prompt, f"Create a presentation on: {user_message}. Output plain JSON only.")
            
        try:
            json_text = res.get("response", "[]")
            # Extra text ko bypass karke sirf JSON array nikalna
            match = re.search(r'\[\s*\{.*?\}\s*\]', json_text, re.DOTALL)
            if match:
                json_text = match.group(0)
            else:
                json_text = "[]"
                
            slides_data = json.loads(json_text)
            
            prs = Presentation()
            for slide_data in slides_data:
                slide_layout = prs.slide_layouts[1] # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_data.get("title", "Untitled Slide")
                
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                for point in slide_data.get("content", []):
                    p = tf.add_paragraph()
                    p.text = str(point)
                    
            os.makedirs("temp", exist_ok=True)
            filename = f"presentation_{int(time.time())}.pptx"
            filepath = os.path.join("temp", filename)
            prs.save(filepath)
            
            return {
                "response": f"Lo bhai! Tumhari PPT ban gayi: <a href='/temp/{filename}' download style='color:#00ffa3;text-decoration:underline'>Download PPT</a>",
                "provider": provider.name,
                "model": provider.model_id
            }
        except Exception as e:
            return {"response": f"PPT generate karne mein error aaya: {str(e)}\n\nAI Output:\n{res.get('response')}"}

    if provider.name == "Google Gemini":
        return await _call_gemini(provider, user_message, full_message)
    else:
        return await _call_openai_compatible(provider, user_message, full_message)


async def _call_gemini(provider: ProviderConfig, system_prompt: str, message: str) -> dict:
    url = f"{provider.base_url}/models/{provider.model_id}:generateContent?key={provider.api_key}"
    payload = {
        "system_instruction": {"parts": [{"text": system_prompt if system_prompt else "You are Manthan AI, a helpful coding assistant."}]},
        "contents": [{"role": "user", "parts": [{"text": message}]}],
        "generationConfig": {"temperature": 0.3, "maxOutputTokens": 4096}
    }
    try:
        async with httpx.AsyncClient(timeout=45) as client:
            resp = await client.post(url, json=payload)
            if resp.status_code != 200:
                return {"response": f"Gemini Error ({resp.status_code}): {resp.text}"}
            data = resp.json()
        text = data["candidates"][0]["content"]["parts"][0]["text"]
        return {"response": text, "provider": provider.name, "model": provider.model_id}
    except Exception as e:
        return {"response": f"Gemini Error: {str(e)}"}


async def _call_gemini_vision(provider: ProviderConfig, message: str, file_bytes: bytes, mime_type: str) -> dict:
    url = f"{provider.base_url}/models/{provider.model_id}:generateContent?key={provider.api_key}"
    base64_file = base64.b64encode(file_bytes).decode("utf-8")
    payload = {
        "contents": [{
            "parts": [
                {"text": message or "Describe this image."},
                {"inline_data": {"mime_type": mime_type, "data": base64_file}}
            ]
        }]
    }
    try:
        async with httpx.AsyncClient(timeout=45) as client:
            resp = await client.post(url, json=payload)
            if resp.status_code != 200:
                return {"response": f"Gemini Vision Error ({resp.status_code}): {resp.text}"}
            data = resp.json()
        text = data["candidates"][0]["content"]["parts"][0]["text"]
        return {"response": text, "provider": "Google Gemini Vision", "model": provider.model_id}
    except Exception as e:
        return {"response": f"Gemini Vision Error: {str(e)}"}


async def _call_openai_compatible(provider: ProviderConfig, system_prompt: str, message: str) -> dict:
    url = f"{provider.base_url}/chat/completions"
    headers = {"Authorization": f"Bearer {provider.api_key}", "Content-Type": "application/json"}
    payload = {
        "model": provider.model_id,
        "messages": [
            {"role": "system", "content": system_prompt if system_prompt else "You are Manthan AI, a helpful coding assistant."},
            {"role": "user",   "content": message}
        ],
        "temperature": 0.3,
        "max_tokens": 4096,
    }
    try:
        async with httpx.AsyncClient(timeout=45) as client:
            resp = await client.post(url, headers=headers, json=payload)
            if resp.status_code != 200:
                return {"response": f"API Error ({resp.status_code}): {resp.text}"}
            data = resp.json()
        return {"response": data["choices"][0]["message"]["content"], "provider": provider.name, "model": provider.model_id}
    except Exception as e:
        return {"response": f"Provider Error: {str(e)}"}